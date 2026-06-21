import sys
from pathlib import Path
import io

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Corporate Color Palette
COLOR_PRIMARY = RGBColor(26, 115, 232)      # Indigo Blue
COLOR_SECONDARY = RGBColor(0, 137, 123)    # Teal
COLOR_DARK_TEXT = RGBColor(32, 33, 36)     # Charcoal Dark
COLOR_LIGHT_TEXT = RGBColor(95, 99, 104)   # Slate Gray
COLOR_LIGHT_BG = RGBColor(248, 249, 250)   # Off-white
COLOR_WHITE = RGBColor(255, 255, 255)


def add_title_slide(prs, customer_name, product_name):
    """Creates a premium title slide with a split color layout."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Left colored accent panel
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(4.5), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLOR_PRIMARY
    left_panel.line.fill.background() # No border
    
    # Left Panel Brand Text (Shodhana AI)
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(3.5), Inches(1))
    tf = brand_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SHODHANA AI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Outfit'
    
    p2 = tf.add_paragraph()
    p2.text = "Growth Engine & Sourcing Pitch"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = 'Outfit'
    
    # Right Main Content Title
    title_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.2), Inches(7.8), Inches(4))
    tf_main = title_box.text_frame
    tf_main.word_wrap = True
    
    p_title = tf_main.paragraphs[0]
    p_title.text = f"Duloxetine Supply Opportunity"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY
    p_title.font.name = 'Outfit'
    
    p_sub = tf_main.add_paragraph()
    p_sub.text = f"Custom Sourcing Proposal for {customer_name}"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_DARK_TEXT
    p_sub.font.name = 'Outfit'
    p_sub.space_before = Pt(12)
    
    p_product = tf_main.add_paragraph()
    p_product.text = f"Target Product Area: {product_name}"
    p_product.font.size = Pt(16)
    p_product.font.color.rgb = COLOR_SECONDARY
    p_product.font.name = 'Roboto'
    p_product.space_before = Pt(20)


def create_header(slide, title_text):
    """Adds a standard header band to content slides."""
    # Top banner bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLOR_PRIMARY
    header_bar.line.fill.background()
    
    # Header Text
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Outfit'


def add_bullets_slide(prs, title, bullets, speaker_note=""):
    """Standard content slide with bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    create_header(slide, title)
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.font.name = 'Roboto'
        p.space_after = Pt(14)
        p.level = 0
        
    if speaker_note and slide.notes_slide:
        slide.notes_slide.notes_text_frame.text = speaker_note


def add_profile_slide(prs, customer_name, country, qty_str, val_str, price_str, shipments, first_ship, last_ship, observed_supplier):
    """Creates a split 2-column detail slide for the customer profile."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    create_header(slide, f"Sourcing Profile: {customer_name}")
    
    # Left Column: Key Volume & Value Metrics
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_LIGHT_BG
    left_box.line.color.rgb = COLOR_PRIMARY
    
    tb_left = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.6))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p_l_title = tf_l.paragraphs[0]
    p_l_title.text = "Sourcing Scale & Volumes"
    p_l_title.font.size = Pt(20)
    p_l_title.font.bold = True
    p_l_title.font.color.rgb = COLOR_PRIMARY
    p_l_title.font.name = 'Outfit'
    p_l_title.space_after = Pt(16)
    
    metrics = [
        f"• Destination Market: {country}",
        f"• Total Sourced Volume: {qty_str}",
        f"• Total Estimated Value: {val_str}",
        f"• Average Sourced Price/KG: {price_str}",
        f"• Total Shipments Recorded: {shipments}",
    ]
    for m in metrics:
        p = tf_l.add_paragraph()
        p.text = m
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.font.name = 'Roboto'
        p.space_after = Pt(10)
        
    # Right Column: Supply Timeline & Observed Exporter
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_LIGHT_BG
    right_box.line.color.rgb = COLOR_SECONDARY
    
    tb_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.6))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    
    p_r_title = tf_r.paragraphs[0]
    p_r_title.text = "Observed Supply Pattern"
    p_r_title.font.size = Pt(20)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = COLOR_SECONDARY
    p_r_title.font.name = 'Outfit'
    p_r_title.space_after = Pt(16)
    
    patterns = [
        f"• Observed Current Exporter: {observed_supplier}",
        f"• First Sourced Shipment: {first_ship}",
        f"• Latest Sourced Shipment: {last_ship}",
        "",
        "Note: Observed average price and shipment cadence suggest a recurring supply relationship. Sourcing patterns should be validated by the BD team during introductory outreach.",
    ]
    for i, pat in enumerate(patterns):
        p = tf_r.add_paragraph()
        p.text = pat
        p.font.size = Pt(14) if "Note:" in pat else Pt(16)
        p.font.color.rgb = COLOR_LIGHT_TEXT if "Note:" in pat else COLOR_DARK_TEXT
        p.font.italic = ("Note:" in pat)
        p.font.name = 'Roboto'
        p.space_after = Pt(10)


def generate_pitch_pptx(opportunity):
    """Generates an editable .pptx in-memory buffer."""
    if not PPTX_AVAILABLE:
        # Graceful fallback: return a simple text explanation
        buffer = io.BytesIO()
        text_data = f"PowerPoint generation not available. Customer: {opportunity.get('importer', 'N/A')}\n"
        buffer.write(text_data.encode('utf-8'))
        buffer.seek(0)
        return buffer, False

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    opp = opportunity.get("opportunity", opportunity)
    customer_name = opp.get("importer") or opp.get("raw_company_name") or "Target Client"
    product_name = opp.get("product") or "Duloxetine API / Pellets"
    country = opp.get("country") or "Global"
    
    # Qty, Val, Price Formatting
    qty_val = opp.get("total_quantity_kg") or 0
    qty_str = f"{qty_val:,.2f} KG" if qty_val else "Not available"
    
    val_val = opp.get("total_value_usd") or 0
    val_str = f"${val_val:,.2f}" if val_val else "Not available"
    
    price_val = opp.get("avg_price_per_kg") or 0
    price_str = f"${price_val:,.2f}/KG" if price_val else "Not available"
    
    shipments = str(opp.get("shipment_count") or "N/A")
    first_ship = opp.get("first_shipment_date") or "N/A"
    last_ship = opp.get("last_shipment_date") or "N/A"
    observed_supplier = opp.get("current_supplier") or opp.get("exporter") or "Competitor Supplier"
    shodhana_status = opp.get("shodhana_status") or "New Account Sourcing"
    score = str(opp.get("score") or "N/A")
    
    # Slide 1: Title Slide
    add_title_slide(prs, customer_name, product_name)
    
    # Slide 2: Sourcing Profile
    add_profile_slide(
        prs, customer_name, country, qty_str, val_str, price_str, 
        shipments, first_ship, last_ship, observed_supplier
    )
    
    # Slide 3: Sourcing Landscape
    landscape_bullets = [
        f"Customer observed supplier context: {observed_supplier}.",
        f"Account Sourcing Category: {shodhana_status}.",
        f"Opportunity Priority Score: {score}/100.",
        "Observed shipments suggest stable, recurring buying behavior.",
        "Targeting objective: Position Shodhana as an alternative, highly dependable source of supply to mitigate buyer's vendor risk."
    ]
    add_bullets_slide(prs, "Supply Landscape & Competitors", landscape_bullets, 
                      f"Review competitor {observed_supplier} strengths before outreach.")
    
    # Slide 4: Commercial & Positioning Strategy
    if "PELLET" in product_name.upper():
        cap_note = "Position Shodhana's semi-formulation/pellets manufacturing strengths and custom release profiles."
    else:
        cap_note = "Position Shodhana's Duloxetine API DMF status, impurity profiles, and high manufacturing capacity."
        
    price_avg_m = opp.get("market_avg_price_per_kg") or 0
    price_diff = opp.get("price_difference") or 0
    if price_avg_m and price_val:
        if price_diff > 0:
            price_note = f"Customer average price (${price_val:,.2f}) appears above market average (${price_avg_m:,.2f}). Position on competitive supply terms without triggering a price war."
        else:
            price_note = f"Customer average price (${price_val:,.2f}) is close to or below market average (${price_avg_m:,.2f}). Position on quality, DMF support, and supply reliability rather than price."
    else:
        price_note = "Commercial positioning should lead with DMF status, regulatory certifications, and reliable delivery terms."
        
    comm_bullets = [
        cap_note,
        price_note,
        "Flexible commercial models: Support DMF/non-DMF commercial setups based on target country regulations.",
        "Supply security: Offer long-term supply agreements with guaranteed volume allocations to prevent supply-chain disruption.",
    ]
    add_bullets_slide(prs, "Commercial Sourcing Strategy", comm_bullets, 
                      "Frame pricing as an open commercial discussion following validation of specification.")
    
    # Slide 5: Strategic Outreach Timeline
    outreach_bullets = [
        "Day 0: Introductory outreach - Send focused Shodhana capability deck and Duloxetine certifications.",
        "Day 3: Capability follow-up - Share regulatory credentials (DMF status) and check validation guidelines.",
        "Day 7: Technical alignment call - Schedule short discussion to align on specifications, mesh sizes, or impurity limits.",
        "Day 14: Sourcing proposal - Prepare a formal price/volume proposal reviewed by the business team.",
        "Day 30: Account review - Track validation progress or re-prioritize based on shipment window cycles."
    ]
    add_bullets_slide(prs, "Outreach Sourcing Timeline", outreach_bullets, 
                      "Sales team should customize outreach frequency based on response rates.")
    
    # Slide 6: Sourcing Proposal & Next Steps
    next_steps = [
        "1. Confirm technical product specification (e.g. mesh size, assay limits, impurity profile).",
        "2. Identify customer's validation timeline and documentation requirements (DMF/CEP/Technical pack).",
        "3. Understand expected annual/quarterly volume allocations.",
        "4. Schedule technical introductory call with customer's sourcing/QA team.",
        "5. Prepare commercial proposal and volume reservation quote upon specification confirmation."
    ]
    add_bullets_slide(prs, "Sourcing Collaboration & Next Steps", next_steps, 
                      "Ensure regulatory team is briefed on target market requirements before the technical call.")

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer, True
